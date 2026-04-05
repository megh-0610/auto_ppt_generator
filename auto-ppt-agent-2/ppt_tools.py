from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from image_services import download_image
import os


OUTPUT_FILE = "outputs/generated.pptx"


# COLORS (from your image)
BG_COLOR = RGBColor(240, 240, 240)
HEADER_COLOR = RGBColor(200, 200, 200)
TEXT_COLOR = RGBColor(50, 50, 50)


# ------------------ EXPAND CONTENT ------------------
def format_bullets(points):

    expanded = []

    for p in points:
        p = str(p).strip()

        if len(p) > 5:
            expanded.append(
                p + " — explained clearly"
            )

    return expanded[:5]


# ------------------ PPT CREATION ------------------
def create_ppt(slides):

    prs = Presentation()

    os.makedirs("outputs", exist_ok=True)
    os.makedirs("temp_images", exist_ok=True)

    for i, s in enumerate(slides):

        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # ================= BACKGROUND =================
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(7.5)
        )

        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_COLOR
        bg.line.fill.background()

        slide.shapes._spTree.remove(bg._element)
        slide.shapes._spTree.insert(2, bg._element)

        # ================= HEADER =================
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(1)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = HEADER_COLOR
        header.line.fill.background()

        # ================= TITLE =================
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.2),
            Inches(9), Inches(0.6)
        )

        title_box.fill.background()
        title_box.line.fill.background()

        tf = title_box.text_frame
        tf.text = s["title"]

        p = tf.paragraphs[0]
        p.font.size = Pt(34)
        p.font.bold = True
        p.font.name = "Cavolini"
        p.font.color.rgb = TEXT_COLOR
        p.alignment = PP_ALIGN.CENTER

        # ================= CONTENT =================
        bullets = format_bullets(s["points"])

        content_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(1.8),
            Inches(7), Inches(4)
        )

        content_box.fill.background()
        content_box.line.fill.background()

        tf = content_box.text_frame
        tf.word_wrap = True

        if bullets:
            tf.text = bullets[0]

        for b in bullets[1:]:
            para = tf.add_paragraph()
            para.text = b
            para.level = 0

        # styling
        for para in tf.paragraphs:
            para.font.size = Pt(22)
            para.font.name = "Angsana New"
            para.font.color.rgb = TEXT_COLOR
            para.alignment = PP_ALIGN.LEFT
            para.space_after = Pt(10)

        # ================= IMAGE =================
        img_path = f"temp_images/image_{i}.png"

        result = download_image(s["image"], img_path)

        if result and os.path.exists(img_path):
            slide.shapes.add_picture(
                img_path,
                Inches(6.5),
                Inches(1.5),
                Inches(3.2),
                Inches(4)
            )

    prs.save(OUTPUT_FILE)

    return OUTPUT_FILE